import os
import ffmpeg
import whisper
import yt_dlp
from datetime import datetime
from pathlib import Path
from PyPDF2 import PdfReader
from bs4 import BeautifulSoup
import requests
from groq import Groq
from utils import split_audio_file_into_parts
from time import sleep
from docx import Document
from pptx import Presentation

groq_api_key = os.environ.get("groq_api_keys")

def transcribe_audio_files(files, urls, language, model_size, USE_GROQ):
    transcript = ""
    temp_dir = './temp'
    os.makedirs(temp_dir, exist_ok=True)
    text_files = []
    audio_files = []
    pdf_files = []
    docx_files = []
    pptx_files = []

    print("Список файлов:")
    for file in files:
        print(file)
        extension = os.path.splitext(file)[-1].lower()
        if extension in ['.txt', '.srt']:
            text_files.append(file)
        elif extension == '.pdf':
            pdf_files.append(file)
        elif extension in ['.docx']:
            docx_files.append(file)
        elif extension in ['.pptx']:
            pptx_files.append(file)
        elif extension in ['.mp3', '.wav', '.mp4']:
            audio_files.append(file)
        else:
            print(f"Неподдерживаемый формат файла: {extension}")

    if text_files:
        combined_text = ""
        for text_file in text_files:
            with open(text_file, "r", encoding="utf-8") as f:
                combined_text += f.read() + "\n"
        transcript = combined_text

    if pdf_files:
        combined_pdf_text = ""
        for pdf_file in pdf_files:
            try:
                reader = PdfReader(pdf_file)
                for page in reader.pages:
                    combined_pdf_text += page.extract_text() + "\n"
            except Exception as e:
                print(f"Ошибка при чтении PDF файла {pdf_file}: {e}")
        transcript += combined_pdf_text

    print("Список ссылок:")
    for url in urls:
        print(url)
        if "youtube.com" in url or "youtu.be" in url:
            print(f"Загрузка YouTube видео с {url}")
            audio_path = download_youtube_audio(url, temp_dir)
            if audio_path:
                audio_files.append(audio_path)
        else:
            print(f"Парсинг текста с {url}")
            page_text = scrape_text_from_url(url)
            if page_text:
                transcript += "\n" + page_text

    if docx_files:
        for word_file in docx_files:
            try:
                transcript += process_word_file(word_file) + "\n"
            except Exception as e:
                print(f"Ошибка при обработке Word файла {word_file}: {e}")

    if pptx_files:
        for pptx_file in pptx_files:
            try:
                transcript += process_pptx_file(pptx_file) + "\n"
            except Exception as e:
                print(f"Ошибка при обработке PowerPoint файла {pptx_file}: {e}")

    if not audio_files:
        print("Аудиофайлы отсутствуют для обработки")
    else:
        audio_transcript = process_audio_files(audio_files, temp_dir, language, model_size, USE_GROQ)
        transcript += "\n" + audio_transcript

    output_dir = './output_transcribed'
    os.makedirs(output_dir, exist_ok=True)
    current_time = datetime.now().strftime("%d-%m-%Y-%H-%M-%S")
    with open(f"{output_dir}/transcript_{current_time}.txt", "w", encoding="utf-8") as f:
        f.write(transcript.strip())

    if os.path.exists(temp_dir):
        for temp_file in os.listdir(temp_dir):
            os.remove(os.path.join(temp_dir, temp_file))
        os.rmdir(temp_dir)

    return transcript

def process_audio_files(audio_files, temp_dir, language, model_size, USE_GROQ):
    audio_paths = []

    for file in audio_files:
        extension = os.path.splitext(file)[-1].lower()
        if extension == '.mp4':
            wav_file = os.path.join(temp_dir, f"{Path(file).stem}.wav")
            ffmpeg.input(file).output(wav_file, format='wav').run()
            audio_paths.append(wav_file)
        elif extension in ['.mp3', '.wav']:
            audio_paths.append(file)

    if len(audio_paths) > 1:
        combined_audio = os.path.join(temp_dir, 'combined.wav')
        inputs = [ffmpeg.input(f) for f in audio_paths]
        ffmpeg.concat(*inputs, v=0, a=1).output(combined_audio, acodec='pcm_s16le').run()
        final_audio = combined_audio
    else:
        final_audio = audio_paths[0]

    mp3_audio = os.path.join(temp_dir, 'converted.mp3')
    ffmpeg.input(final_audio).output(mp3_audio, audio_bitrate='256k').run()

    audio_parts = split_audio_file_into_parts(mp3_audio, max_size_mb=24)

    audio_transcript = ""

    if USE_GROQ:
        print("Использование Whisper GROQ для транскрибирования текста")
        client = Groq(api_key=groq_api_key)

        for idx, audio_file in enumerate(audio_parts):
            with open(audio_file, "rb") as file:
                transcription = client.audio.transcriptions.create(
                    file=(audio_file, file.read()),
                    model="whisper-large-v3-turbo",
                    language=language
                )
                audio_transcript += transcription.text + "\n"
                print(f"Часть {idx+1} транскрибирована.")
                if len(audio_parts) > 1:
                    print("Засыпаем на 65 секунд, из-за ограничений API")
                    sleep(65)
    else:
        print("Использование локального Whisper для транскрибирования текста")
        model = whisper.load_model(model_size)
        result = model.transcribe(audio=mp3_audio, language=language, verbose=False)
        audio_transcript = result.get('text', '')

    return audio_transcript

def scrape_text_from_url(url):
    try:
        response = requests.get(url)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')

        tags_to_parse = ['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'a', 'span', 'ul']
        extracted_text = []
        for tag in tags_to_parse:
            elements = soup.find_all(tag)
            for element in elements:
                text = element.get_text(strip=True)
                if text:
                    extracted_text.append(text)

        return "\n".join(extracted_text)

    except requests.exceptions.RequestException as e:
        print(f"Ошибка при получении URL {url}: {e}")
        return None
    except Exception as e:
        print(f"Неожиданная ошибка при парсинге {url}: {e}")
        return None

def download_youtube_audio(url, temp_dir):
    ydl_opts = {
        'format': 'bestaudio/best',
        'outtmpl': os.path.join(temp_dir, '%(title)s.%(ext)s'),
        'postprocessors': [
            {
                'key': 'FFmpegExtractAudio',
                'preferredcodec': 'mp3',
                'preferredquality': '192',
            }
        ],
        'keepvideo': False,
    }

    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info_dict = ydl.extract_info(url, download=True)
            base_filename_with_ext = ydl.prepare_filename(info_dict)
            mp3_filename = base_filename_with_ext.rsplit('.', 1)[0] + '.mp3'

            if not os.path.exists(mp3_filename):
                print(f"Файл {mp3_filename} не найден после обработки")
                return None

            return mp3_filename
    except Exception as e:
        print(f"Ошибка при загрузке YouTube-видео: {e}")
        return None
    
def process_word_file(file_path):
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return "\n".join(text)

def process_pptx_file(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text.append(paragraph.text)
    return "\n".join(text)